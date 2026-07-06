import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_benchmark_and_hyperparameter_slides(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return

    try:
        prs = pptx.Presentation(pptx_path)
    except Exception as e:
        print(f"Could not open {pptx_path}: {e}")
        return

    print(f"Initial slide count: {len(prs.slides)}")

    # 1. Find Slide 15: LightGBM Hyperparameter Tuning Experiments
    slide15_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and 'LightGBM Hyperparameter Tuning Experiments' in shape.text_frame.text:
                slide15_idx = i
                break
        if slide15_idx is not None:
            break

    if slide15_idx is None:
        print("Could not find LightGBM Hyperparameter Tuning slide, defaulting to index 14.")
        slide15_idx = 14

    print(f"Found LightGBM Hyperparameter slide at index {slide15_idx} (Slide {slide15_idx+1}).")

    # Check if our new Hyperparameter slide is already inserted
    next_idx = slide15_idx + 1
    has_hyp_slide = False
    if next_idx < len(prs.slides):
        for shape in prs.slides[next_idx].shapes:
            if shape.has_text_frame and 'Baseline Model Hyperparameters' in shape.text_frame.text:
                has_hyp_slide = True
                break

    if has_hyp_slide:
        print("Baseline Hyperparameter slide already exists.")
    else:
        # Create Slide A: Baseline Hyperparameter Configurations
        blank_layout = prs.slide_layouts[6] # Blank layout
        slide_a = prs.slides.add_slide(blank_layout)

        # Move slide_a to next_idx (right after LightGBM Hyperparameters)
        sldIdLst = prs.slides._sldIdLst
        slide_el = sldIdLst[-1]
        sldIdLst.remove(slide_el)
        sldIdLst.insert(next_idx, slide_el)

        # Add Title & Subtitle
        t_box = slide_a.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.1))
        t_tf = t_box.text_frame
        t_tf.word_wrap = True
        tp = t_tf.paragraphs[0]
        tp.text = "Baseline Model Hyperparameters & Why LightGBM Wins"
        tp.font.name = "Calibri"
        tp.font.size = Pt(26)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(16, 44, 87) # Navy

        subp = t_tf.add_paragraph()
        subp.text = "Comparative Analysis of Training Parameters across Linear Regression, Random Forest, XGBoost, and LightGBM"
        subp.font.name = "Calibri"
        subp.font.size = Pt(15)
        subp.font.color.rgb = RGBColor(100, 110, 120)

        # Add Table: 5 rows (1 header + 4 data), 3 cols
        rows, cols = 5, 3
        t_left, t_top, t_width, t_height = Inches(0.6), Inches(1.6), Inches(12.1), Inches(3.8)
        t_shape = slide_a.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
        table = t_shape.table

        table.columns[0].width = Inches(2.6) # Model Family
        table.columns[1].width = Inches(4.5) # Key Hyperparameters
        table.columns[2].width = Inches(5.0) # Automotive ECU Feasibility & Explanation

        headers = ["Model Family", "Key Hyperparameter Configuration Used", "Automotive ECU Feasibility & Explanation"]
        data = [
            [
                "Linear Regression\n(Severson Baseline)",
                "• Ridge Regularization: alpha = 1.0\n• Feature Scaling: StandardScaler()\n• Loss Function: L2 Mean Squared Error",
                "DEPLOYABLE BUT HIGH ERROR (> 150c MAE)\nLinear models assume straight-line degradation, completely missing the late-life capacity plunge at 80% SOH."
            ],
            [
                "Random Forest\nRegressor",
                "• n_estimators = 100 trees | max_depth = 12\n• Bootstrap Sampling: True | max_features = 'sqrt'\n• Split Criterion: Squared Error | n_jobs = -1",
                "FAILED: OUT OF FLASH MEMORY (28.1 MB)\nBuilds 100 deep, unpruned trees. Serialized object is 28.1 MB, exceeding the strict 1 MB Flash limit of low-cost ECUs by 28x!"
            ],
            [
                "XGBoost Regressor\n(Gradient Boosting)",
                "• n_estimators = 300 | learning_rate = 0.05\n• max_depth = 6 | subsample = 0.8 | colsample = 0.8\n• Tree Growth: Level-wise (Depth-First)",
                "FAILED: EXCEEDS FLASH LIMIT (1.36 MB)\nLevel-wise tree growth across 300 trees creates excessive node bloat (1.36 MB), exceeding standard 1 MB automotive microcontroller limits."
            ],
            [
                "Our LightGBM\nFramework (Ours)",
                "• n_estimators = 300 | learning_rate = 0.05\n• max_depth = 7 | num_leaves = 31 | min_child = 20\n• Tree Growth: Leaf-wise (Best-First) + GOSS & EFB",
                "OPTIMAL EDGE WINNER (809 KB Flash, 0.41 ms)\nLeaf-wise pruning and histogram binning compress the model into 809 KB, achieving the lowest error (79.91c MAE) within sub-ms execution!"
            ]
        ]

        for col_i, h_text in enumerate(headers):
            cell = table.cell(0, col_i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(16, 44, 87)
            p = cell.text_frame.paragraphs[0]
            p.text = h_text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Calibri"
            p.font.size = Pt(13.5)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

        for row_i, row_data in enumerate(data):
            fill_col = RGBColor(245, 247, 250) if row_i % 2 == 0 else RGBColor(255, 255, 255)
            if row_i == 3: # Highlight LightGBM
                fill_col = RGBColor(235, 245, 240) # Soft green tint

            for col_i, val in enumerate(row_data):
                cell = table.cell(row_i + 1, col_i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_col
                tf = cell.text_frame
                tf.word_wrap = True
                tf.margin_top = Inches(0.08)
                tf.margin_bottom = Inches(0.08)
                tf.margin_left = Inches(0.12)
                tf.margin_right = Inches(0.12)

                p = tf.paragraphs[0]
                lines = val.split('\n')
                p.text = lines[0]
                p.font.name = "Calibri"
                p.font.size = Pt(12)
                p.font.bold = (col_i == 0 or row_i == 3 or "FAILED" in lines[0] or "DEPLOYABLE" in lines[0] or "OPTIMAL" in lines[0])
                if "FAILED" in lines[0] or "HIGH ERROR" in lines[0]:
                    p.font.color.rgb = RGBColor(180, 40, 40)
                elif "OPTIMAL" in lines[0]:
                    p.font.color.rgb = RGBColor(25, 135, 84)
                else:
                    p.font.color.rgb = RGBColor(33, 37, 41)

                for line in lines[1:]:
                    p2 = tf.add_paragraph()
                    p2.text = line
                    p2.font.name = "Calibri"
                    p2.font.size = Pt(11.5)
                    p2.font.color.rgb = RGBColor(50, 60, 70)
                    if col_i == 0:
                        p2.font.bold = True

        # Takeaway Box below table
        box = slide_a.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t_left, Inches(5.65), t_width, Inches(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 244, 248)
        box.line.color.rgb = RGBColor(16, 44, 87)
        box.line.width = Pt(1.5)

        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.2)
        btf.margin_top = Inches(0.1)

        bp0 = btf.paragraphs[0]
        bp0.text = "Why Hyperparameter Design is Critical for Automotive Microcontrollers:"
        bp0.font.name = "Calibri"
        bp0.font.size = Pt(13.5)
        bp0.font.bold = True
        bp0.font.color.rgb = RGBColor(16, 44, 87)

        bp1 = btf.add_paragraph()
        bp1.text = "• Level-Wise vs. Leaf-Wise Tree Growth: Random Forest and XGBoost grow full depth trees across all nodes, causing massive Flash memory bloat (>1.3 MB to 28 MB) that crashes standard 1 MB automotive ECUs.\n• LightGBM's Efficiency Secret: By combining Leaf-Wise pruning (num_leaves=31) with Gradient-Based One-Side Sampling (GOSS) and Exclusive Feature Bundling (EFB), LightGBM achieves the highest predictive accuracy while using only 809 KB Flash!"
        bp1.font.name = "Calibri"
        bp1.font.size = Pt(12)
        bp1.font.color.rgb = RGBColor(33, 37, 41)

        print("Successfully inserted Slide A: Baseline Hyperparameter Configurations after Slide 15.")

    # 2. Find Unseen Test Cohort Validation Performance slide (which is now after our new slide)
    slide_unseen_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and 'Unseen Test Cohort Validation Performance' in shape.text_frame.text:
                slide_unseen_idx = i
                break
        if slide_unseen_idx is not None:
            break

    if slide_unseen_idx is None:
        print("Could not find Unseen Test Cohort slide, defaulting to index 16.")
        slide_unseen_idx = 16

    print(f"Found Unseen Test Cohort slide at index {slide_unseen_idx} (Slide {slide_unseen_idx+1}).")

    # Check if Slide B is already inserted after Unseen Test Cohort
    next_idx_b = slide_unseen_idx + 1
    has_bench_slide = False
    if next_idx_b < len(prs.slides):
        for shape in prs.slides[next_idx_b].shapes:
            if shape.has_text_frame and 'Empirical Benchmarking: LightGBM vs Baseline Models' in shape.text_frame.text:
                has_bench_slide = True
                break

    if has_bench_slide:
        print("Empirical Benchmarking slide already exists.")
    else:
        # Create Slide B: Empirical Benchmarking Comparison (Table 6 & Fig. 6)
        blank_layout = prs.slide_layouts[6]
        slide_b = prs.slides.add_slide(blank_layout)

        # Move slide_b to next_idx_b
        sldIdLst = prs.slides._sldIdLst
        slide_el = sldIdLst[-1]
        sldIdLst.remove(slide_el)
        sldIdLst.insert(next_idx_b, slide_el)

        # Add Title & Subtitle
        t_box = slide_b.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.1))
        t_tf = t_box.text_frame
        t_tf.word_wrap = True
        tp = t_tf.paragraphs[0]
        tp.text = "Empirical Benchmarking: LightGBM vs Baseline Models"
        tp.font.name = "Calibri"
        tp.font.size = Pt(26)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(16, 44, 87)

        subp = t_tf.add_paragraph()
        subp.text = "Evaluating Predictive Accuracy (MAE & R²), Flash Size, MCU Loop Latency, and Automotive ECU Feasibility (Table 6 & Fig. 6)"
        subp.font.name = "Calibri"
        subp.font.size = Pt(15)
        subp.font.color.rgb = RGBColor(100, 110, 120)

        # Add Table 6 at the top (Top 1.55, Left 0.6, Width 12.1, Height 2.1)
        rows, cols = 5, 6
        t_left, t_top, t_width, t_height = Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.1)
        t_shape = slide_b.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
        table = t_shape.table

        table.columns[0].width = Inches(2.3) # Model
        table.columns[1].width = Inches(1.5) # MAE
        table.columns[2].width = Inches(1.3) # R2
        table.columns[3].width = Inches(1.4) # Flash Size
        table.columns[4].width = Inches(1.4) # MCU Loop
        table.columns[5].width = Inches(4.2) # Automotive ECU Feasibility

        headers = ["Model", "MAE (Cycles)", "R² (%)", "Flash Size", "MCU Loop", "Automotive ECU Feasibility"]
        data = [
            ["Linear Regression", "152.90", "56.12%", "1.2 KB", "0.22 ms", "Deployable but high error (> 150c)."],
            ["Random Forest", "80.18", "79.15%", "28.1 MB", "1.85 ms", "FAILED: Exceeds 1 MB Flash limit."],
            ["XGBoost Regressor", "82.53", "80.25%", "1.36 MB", "0.68 ms", "FAILED: Exceeds 1 MB Flash limit."],
            ["LightGBM (Ours)", "79.91", "81.62%", "809 KB", "0.41 ms", "OPTIMAL: Fits Flash, < 1.5 ms loop."]
        ]

        for col_i, h_text in enumerate(headers):
            cell = table.cell(0, col_i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(16, 44, 87)
            p = cell.text_frame.paragraphs[0]
            p.text = h_text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Calibri"
            p.font.size = Pt(12.5)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

        for row_i, row_data in enumerate(data):
            fill_col = RGBColor(245, 247, 250) if row_i % 2 == 0 else RGBColor(255, 255, 255)
            if row_i == 3: # LightGBM
                fill_col = RGBColor(235, 245, 240)

            for col_i, val in enumerate(row_data):
                cell = table.cell(row_i + 1, col_i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_col
                tf = cell.text_frame
                tf.word_wrap = True
                tf.margin_top = Inches(0.05)
                tf.margin_bottom = Inches(0.05)
                tf.margin_left = Inches(0.08)
                tf.margin_right = Inches(0.08)

                p = tf.paragraphs[0]
                p.text = val
                p.alignment = PP_ALIGN.CENTER if col_i in [1, 2, 3, 4] else PP_ALIGN.LEFT
                p.font.name = "Calibri"
                p.font.size = Pt(11.5)
                p.font.bold = (col_i == 0 or row_i == 3 or "FAILED" in val or "OPTIMAL" in val)
                if "FAILED" in val:
                    p.font.color.rgb = RGBColor(180, 40, 40)
                elif "OPTIMAL" in val or (row_i == 3 and col_i in [1, 2]):
                    p.font.color.rgb = RGBColor(25, 135, 84)
                else:
                    p.font.color.rgb = RGBColor(33, 37, 41)

        # Below Table 6: Left side insert Fig 6 image, Right side insert explanation box
        img_path = r"d:\chandru project\RUL prediction\results\benchmarks\model_comparison_bar_chart.png"
        if os.path.exists(img_path):
            slide_b.shapes.add_picture(img_path, Inches(0.6), Inches(3.85), Inches(6.0), Inches(3.2))
        else:
            print(f"Warning: Image not found at {img_path}")

        # Right side Explanation Card (Left 6.8, Top 3.85, Width 5.9, Height 3.2)
        box = slide_b.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.85), Inches(5.9), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 244, 248)
        box.line.color.rgb = RGBColor(16, 44, 87)
        box.line.width = Pt(1.5)

        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.2)
        btf.margin_right = Inches(0.2)
        btf.margin_top = Inches(0.12)

        bp0 = btf.paragraphs[0]
        bp0.text = "Why LightGBM Outperforms All Baselines:"
        bp0.font.name = "Calibri"
        bp0.font.size = Pt(14)
        bp0.font.bold = True
        bp0.font.color.rgb = RGBColor(16, 44, 87)

        points = [
            "1. Superior Predictive Accuracy (79.91c MAE | 81.62% R²):\nLightGBM reduces error by nearly half compared to Linear Regression (152.90c MAE), successfully capturing the steep late-life capacity plunge that linear models miss entirely.",
            "2. Overcoming Automotive Flash Limits (< 1 MB):\nBoth Random Forest (28.1 MB) and XGBoost (1.36 MB) FAIL automotive feasibility because they exceed the strict 1 MB Flash storage limit of standard ARM Cortex / Infineon microcontrollers.",
            "3. Ultra-Fast Execution (0.41 ms MCU Loop):\nBy evaluating simple integer IF/ELSE thresholds via histogram binning, LightGBM executes in 0.41 ms, leaving 95% of the controller's budget free for motor control & braking!"
        ]

        for pt in points:
            p = btf.add_paragraph()
            lines = pt.split('\n')
            p.text = lines[0]
            p.font.name = "Calibri"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(16, 44, 87) if "1." in lines[0] or "3." in lines[0] else RGBColor(180, 40, 40)
            if "3." in lines[0] or "1." in lines[0]:
                p.font.color.rgb = RGBColor(25, 135, 84)

            p2 = btf.add_paragraph()
            p2.text = lines[1]
            p2.font.name = "Calibri"
            p2.font.size = Pt(11)
            p2.font.color.rgb = RGBColor(33, 37, 41)

        print("Successfully inserted Slide B: Empirical Benchmarking after Unseen Test Cohort slide.")

    try:
        prs.save(pptx_path)
        print(f"Presentation saved successfully to: {pptx_path}")
        print(f"Final slide count: {len(prs.slides)}")
    except PermissionError:
        print(f"PERMISSION ERROR: Please close PowerPoint! File is locked: {pptx_path}")

if __name__ == "__main__":
    add_benchmark_and_hyperparameter_slides(r"D:\chandru project\RUL prediction\reports\presentations\rul final.pptx")
    # Also copy to d:\chandru project\rul final.pptx if desired
    if os.path.exists(r"D:\chandru project\RUL prediction\reports\presentations\rul final.pptx"):
        import shutil
        shutil.copy2(r"D:\chandru project\RUL prediction\reports\presentations\rul final.pptx", r"D:\chandru project\rul final.pptx")
        print("Copied updated deck to D:\\chandru project\\rul final.pptx")
