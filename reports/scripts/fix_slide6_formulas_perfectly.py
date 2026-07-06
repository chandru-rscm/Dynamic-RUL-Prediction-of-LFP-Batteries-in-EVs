import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def fix_formulas_with_xml_baseline(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return
    
    print(f"Opening presentation: {pptx_path}")
    prs = pptx.Presentation(pptx_path)

    # -------------------------------------------------------------------------
    # SLIDE 6 (INDEX 5): REBUILD TABLE 1 IN USER'S EXACT REQUESTED ORDER
    # -------------------------------------------------------------------------
    print("Step 1: Rebuilding Table 1 on Slide 6 in exact requested order...")
    slide6 = prs.slides[5]  # Index 5 is Slide 6
    
    # Remove all shapes on slide 6 except the title (Shape 0)
    for sp in list(slide6.shapes)[1:]:
        sp._element.getparent().remove(sp._element)
        
    # Table dimensions: 9 rows (1 header + 8 features), 4 cols
    rows, cols = 9, 4
    t_left, t_top = Inches(0.5), Inches(1.5)
    t_width, t_height = Inches(12.33), Inches(5.5)
    
    table_shape = slide6.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.25)  # Symbol
    table.columns[1].width = Inches(2.15)  # Feature Name
    table.columns[2].width = Inches(4.93)  # Physical Meaning & Degradation Link
    table.columns[3].width = Inches(4.00)  # Computation Formula
    
    headers = ["Symbol", "Feature Name", "Physical Meaning & Degradation Link", "Computation Formula"]
    
    # Helper to style cells and add formatted runs with direct XML baseline manipulation
    def set_cell_content(cell, runs_list, align=PP_ALIGN.LEFT, is_header=False, row_idx=0):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if is_header:
            cell.fill.fore_color.rgb = RGBColor(16, 44, 87) # Navy
        else:
            cell.fill.fore_color.rgb = RGBColor(245, 247, 250) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            
        p = cell.text_frame.paragraphs[0]
        p.text = ""
        p.alignment = align
        
        for run_data in runs_list:
            text = run_data[0]
            italic = run_data[1]
            bold = run_data[2]
            sub = run_data[3] if len(run_data) > 3 else False
            sup = run_data[4] if len(run_data) > 4 else False
            font_name = run_data[5] if len(run_data) > 5 else ("Cambria Math" if (italic or sub or sup) else "Calibri")
            size_pt = run_data[6] if len(run_data) > 6 else (13.5 if is_header else (9.5 if (sub or sup) else 11.5))
            
            r = p.add_run()
            r.text = text
            r.font.name = font_name
            r.font.size = Pt(size_pt)
            r.font.italic = italic
            r.font.bold = bold
            r.font.color.rgb = RGBColor(255, 255, 255) if is_header else RGBColor(33, 37, 41)
            
            # CRITICAL: Directly set XML baseline attribute for genuine PowerPoint subscript/superscript
            if sub:
                r.font._element.set("baseline", "-25000") # Lower by 25%
            elif sup:
                r.font._element.set("baseline", "30000")  # Raise by 30%

    # Style Header Row
    for col_i, h_text in enumerate(headers):
        set_cell_content(table.cell(0, col_i), [(h_text, False, True, False, False, "Calibri", 13.5)], align=PP_ALIGN.CENTER, is_header=True)
        
    # Feature Data in exact requested order:
    # 1. cycle
    # 2. soh
    # 3. IR
    # 4. Tmean
    # 5. Capacity fade window
    # 6. delta q log variance
    # 7. delta q min
    # 8. delta q mean
    features_data = [
        # 1. cycle (k)
        (
            [("k", True, True, False, False, "Cambria Math", 12)],
            [("Cycle Number", False, True, False, False, "Calibri", 11.5)],
            [("Tracks baseline temporal aging progression and cumulative operational charge-discharge duration of the cell.", False, False, False, False, "Calibri", 11)],
            [
                ("Current operational cycle index ", False, False, False, False, "Calibri", 11.5),
                ("k", True, True, False, False, "Cambria Math", 12),
                (" at checkpoint", False, False, False, False, "Calibri", 11.5)
            ]
        ),
        # 2. SOH
        (
            [("SOH", False, True, False, False, "Calibri", 12)],
            [("State of Health", False, True, False, False, "Calibri", 11.5)],
            [("Normalized remaining capacity ratio; primary macro-level health indicator across cell form factors.", False, False, False, False, "Calibri", 11)],
            [
                ("Q", True, True, False, False, "Cambria Math", 12),
                ("dis", False, False, True, False, "Calibri", 9.5),
                ("(", False, False, False, False, "Cambria Math", 12),
                ("k", True, True, False, False, "Cambria Math", 12),
                (") / Q", True, True, False, False, "Cambria Math", 12),
                ("nominal", False, False, True, False, "Calibri", 9.5)
            ]
        ),
        # 3. IR
        (
            [("IR", False, True, False, False, "Calibri", 12)],
            [("Internal Resistance", False, True, False, False, "Calibri", 11.5)],
            [("Measures ohmic + SEI conduction resistance; primary indicator of electrolyte decomposition and interfacial thickening.", False, False, False, False, "Calibri", 11)],
            [
                ("Mean( ", False, False, False, False, "Cambria Math", 11.5),
                ("dV", True, True, False, False, "Cambria Math", 11.5),
                (" / ", False, False, False, False, "Cambria Math", 11.5),
                ("dI", True, True, False, False, "Cambria Math", 11.5),
                (" )", False, False, False, False, "Cambria Math", 11.5),
                ("10% SOC pulse", False, False, True, False, "Calibri", 9.5)
            ]
        ),
        # 4. Tmean
        (
            [
                ("T", True, True, False, False, "Cambria Math", 12),
                ("mean", False, False, True, False, "Calibri", 9.5)
            ],
            [("Mean Discharge Temp", False, True, False, False, "Calibri", 11.5)],
            [("Monitors thermal stress during operational discharge; drives Arrhenius aging acceleration and SEI growth.", False, False, False, False, "Calibri", 11)],
            [
                ("T̄", True, True, False, False, "Cambria Math", 12),
                ("surface/core", False, False, True, False, "Calibri", 9.5),
                (" over ", False, False, False, False, "Calibri", 11.5),
                ("W", True, True, False, False, "Cambria Math", 12),
                (" = 10", False, False, False, False, "Cambria Math", 12)
            ]
        ),
        # 5. capacity_fade_window (ΔSOH)
        (
            [("ΔSOH", False, True, False, False, "Calibri", 12)],
            [("Capacity Fade Rate", False, True, False, False, "Calibri", 11.5)],
            [("Normalized rate of capacity loss over lookback window; measures local degradation velocity (ΔSOH).", False, False, False, False, "Calibri", 11)],
            [
                ("[ ", False, False, False, False, "Cambria Math", 11.5),
                ("Q", True, True, False, False, "Cambria Math", 11.5),
                ("dis", False, False, True, False, "Calibri", 9.5),
                ("(", False, False, False, False, "Cambria Math", 11.5),
                ("k", True, True, False, False, "Cambria Math", 11.5),
                (" - 10) - ", False, False, False, False, "Cambria Math", 11.5),
                ("Q", True, True, False, False, "Cambria Math", 11.5),
                ("dis", False, False, True, False, "Calibri", 9.5),
                ("(", False, False, False, False, "Cambria Math", 11.5),
                ("k", True, True, False, False, "Cambria Math", 11.5),
                (") ] / ", False, False, False, False, "Cambria Math", 11.5),
                ("Q", True, True, False, False, "Cambria Math", 11.5),
                ("nominal", False, False, True, False, "Calibri", 9.5)
            ]
        ),
        # 6. dQ_log_var
        (
            [
                ("ΔQ", True, True, False, False, "Cambria Math", 12),
                ("log var", False, False, True, False, "Calibri", 9.5)
            ],
            [("IC Curve Log-Variance", False, True, False, False, "Calibri", 11.5)],
            [("Captures thermodynamic peak broadening in incremental capacity curve; acts as an early warning sentinel.", False, False, False, False, "Calibri", 11)],
            [
                ("log[ Var( ", False, False, False, False, "Cambria Math", 11.5),
                ("dQ", True, True, False, False, "Cambria Math", 11.5),
                (" / ", False, False, False, False, "Cambria Math", 11.5),
                ("dV", True, True, False, False, "Cambria Math", 11.5),
                (" ) ]", False, False, False, False, "Cambria Math", 11.5),
                ("W=10", False, False, True, False, "Calibri", 9.5)
            ]
        ),
        # 7. dQ_min
        (
            [
                ("ΔQ", True, True, False, False, "Cambria Math", 12),
                ("min", False, False, True, False, "Calibri", 9.5)
            ],
            [("IC Curve Min Shift", False, True, False, False, "Calibri", 11.5)],
            [("Maximum localized downward shift (valley) in the differential capacity profile over the rolling lookback window.", False, False, False, False, "Calibri", 11)],
            [
                ("min[ Δ( ", False, False, False, False, "Cambria Math", 11.5),
                ("dQ", True, True, False, False, "Cambria Math", 11.5),
                (" / ", False, False, False, False, "Cambria Math", 11.5),
                ("dV", True, True, False, False, "Cambria Math", 11.5),
                (" ) ]", False, False, False, False, "Cambria Math", 11.5),
                ("W=10", False, False, True, False, "Calibri", 9.5)
            ]
        ),
        # 8. dQ_mean
        (
            [
                ("ΔQ", True, True, False, False, "Cambria Math", 12),
                ("mean", False, False, True, False, "Calibri", 9.5)
            ],
            [("IC Curve Mean Shift", False, True, False, False, "Calibri", 11.5)],
            [("Average baseline shift of incremental capacity curve; tracks global thermodynamic drift and stoichiometric imbalance.", False, False, False, False, "Calibri", 11)],
            [
                ("mean[ Δ( ", False, False, False, False, "Cambria Math", 11.5),
                ("dQ", True, True, False, False, "Cambria Math", 11.5),
                (" / ", False, False, False, False, "Cambria Math", 11.5),
                ("dV", True, True, False, False, "Cambria Math", 11.5),
                (" ) ]", False, False, False, False, "Cambria Math", 11.5),
                ("W=10", False, False, True, False, "Calibri", 9.5)
            ]
        ),
    ]
    
    for row_idx, row_data in enumerate(features_data, start=1):
        set_cell_content(table.cell(row_idx, 0), row_data[0], align=PP_ALIGN.CENTER, row_idx=row_idx)
        set_cell_content(table.cell(row_idx, 1), row_data[1], align=PP_ALIGN.LEFT, row_idx=row_idx)
        set_cell_content(table.cell(row_idx, 2), row_data[2], align=PP_ALIGN.LEFT, row_idx=row_idx)
        set_cell_content(table.cell(row_idx, 3), row_data[3], align=PP_ALIGN.LEFT, row_idx=row_idx)
        
    print(" -> Successfully rebuilt Table 1 equations in exact requested order.")

    # -------------------------------------------------------------------------
    # SLIDE 21 (INDEX 20): REFINE F1 SUBSCRIPT IN PERFORMANCE METRICS TABLE
    # -------------------------------------------------------------------------
    print("Step 2: Refining F1 subscript on Slide 21...")
    slide21 = prs.slides[20]
    for sp in slide21.shapes:
        if sp.has_table:
            t = sp.table
            # Row 4 is Overall (F1)
            cell = t.cell(4, 0)
            p = cell.text_frame.paragraphs[0]
            p.text = ""
            p.alignment = PP_ALIGN.CENTER
            
            r1 = p.add_run()
            r1.text = "Overall (F"
            r1.font.name = "Calibri"
            r1.font.size = Pt(11.5)
            r1.font.bold = True
            r1.font.color.rgb = RGBColor(33, 37, 41)
            
            r2 = p.add_run()
            r2.text = "1"
            r2.font.name = "Calibri"
            r2.font.size = Pt(9.5)
            r2.font.bold = True
            r2.font._element.set("baseline", "-25000") # True subscript
            r2.font.color.rgb = RGBColor(33, 37, 41)
            
            r3 = p.add_run()
            r3.text = ")"
            r3.font.name = "Calibri"
            r3.font.size = Pt(11.5)
            r3.font.bold = True
            r3.font.color.rgb = RGBColor(33, 37, 41)
            
            # Also check Formula cell (row 4, col 1)
            cell_f = t.cell(4, 1)
            pf = cell_f.text_frame.paragraphs[0]
            pf.text = ""
            pf.alignment = PP_ALIGN.CENTER
            
            rf1 = pf.add_run()
            rf1.text = "F"
            rf1.font.name = "Cambria Math"
            rf1.font.size = Pt(11.5)
            rf1.font.italic = True
            rf1.font.color.rgb = RGBColor(33, 37, 41)
            
            rf2 = pf.add_run()
            rf2.text = "1"
            rf2.font.name = "Calibri"
            rf2.font.size = Pt(9.5)
            rf2.font._element.set("baseline", "-25000")
            rf2.font.color.rgb = RGBColor(33, 37, 41)
            
            rf3 = pf.add_run()
            rf3.text = " = (2 · Precision · Recall) / (Precision + Recall)  (11)"
            rf3.font.name = "Calibri"
            rf3.font.size = Pt(11.5)
            rf3.font.color.rgb = RGBColor(33, 37, 41)
            break
            
    print(" -> Refined Slide 21 F1 subscripts.")

    # Save presentation
    try:
        prs.save(pptx_path)
        print(f"\nSUCCESS! Saved updated presentation to: {pptx_path}")
        root_path = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(pptx_path))), "rul final.pptx")
        if os.path.exists(root_path):
            import shutil
            shutil.copy2(pptx_path, root_path)
            print(f"Synced copy to root directory: {root_path}")
    except PermissionError:
        print(f"\nPERMISSION ERROR: Please close PowerPoint! File is currently locked: {pptx_path}")

if __name__ == "__main__":
    target = r"D:\chandru project\RUL prediction\reports\presentations\rul final.pptx"
    fix_formulas_with_xml_baseline(target)
