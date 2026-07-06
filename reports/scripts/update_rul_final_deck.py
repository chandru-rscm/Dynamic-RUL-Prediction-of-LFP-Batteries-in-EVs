import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def update_rul_final(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return

    try:
        prs = pptx.Presentation(pptx_path)
    except Exception as e:
        print(f"Could not open {pptx_path}: {e}")
        return

    # 1. Update Confusion Matrix Slide with Formulas
    # Find the Confusion Matrix slide by title
    cm_slide = None
    for slide in prs.slides:
        if slide.shapes and slide.shapes[0].has_text_frame and 'Confusion Matrix' in slide.shapes[0].text_frame.text:
            cm_slide = slide
            break

    if cm_slide:
        # Check if we already added a formula box
        has_formula_box = any(sp.has_text_frame and 'TP + TN' in sp.text_frame.text for sp in cm_slide.shapes)
        if not has_formula_box:
            # Add rounded rectangle box at bottom right
            f_box = cm_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7.2), Inches(4.7), Inches(5.6), Inches(2.3)
            )
            f_box.fill.solid()
            f_box.fill.fore_color.rgb = RGBColor(240, 244, 248) # Soft slate/blue tint
            f_box.line.color.rgb = RGBColor(16, 44, 87) # Navy border
            f_box.line.width = Pt(1.5)

            tf = f_box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.15)
            tf.margin_top = Inches(0.12)

            p0 = tf.paragraphs[0]
            p0.text = "Exact Evaluation Formulas (4,234 Test Samples):"
            p0.font.name = "Calibri"
            p0.font.size = Pt(13)
            p0.font.bold = True
            p0.font.color.rgb = RGBColor(16, 44, 87)

            formulas = [
                "• Accuracy = (TP + TN) / Total = (498 + 3600) / 4234 = 96.79%",
                "• Precision = TP / (TP + FP) = 498 / (498 + 72) = 87.37%",
                "• Recall = TP / (TP + FN) = 498 / (498 + 64) = 88.61%"
            ]

            for f_text in formulas:
                p = tf.add_paragraph()
                p.text = f_text
                p.font.name = "Calibri"
                p.font.size = Pt(12.5)
                p.font.bold = True
                p.font.color.rgb = RGBColor(33, 37, 41)
            print("Successfully added formula card to Confusion Matrix slide.")

    # 2. Add Comparison Slide after Slide 13 (Rolling Lookback Window Table)
    # Find Slide 13 index
    slide13_idx = None
    for i, slide in enumerate(prs.slides):
        if slide.shapes and slide.shapes[0].has_text_frame and 'Rolling Lookback Window' in slide.shapes[0].text_frame.text:
            slide13_idx = i
            break

    if slide13_idx is not None:
        # Check if Slide 14 is already our comparison slide
        next_idx = slide13_idx + 1
        if next_idx < len(prs.slides):
            next_slide = prs.slides[next_idx]
            if next_slide.shapes and next_slide.shapes[0].has_text_frame and 'Architectural Comparison' in next_slide.shapes[0].text_frame.text:
                print("Comparison slide already exists right after Slide 13.")
            else:
                # Create blank slide
                blank_layout = prs.slide_layouts[6] # blank layout
                new_slide = prs.slides.add_slide(blank_layout)

                # Move new slide to next_idx
                sldIdLst = prs.slides._sldIdLst
                slide_el = sldIdLst[-1]
                sldIdLst.remove(slide_el)
                sldIdLst.insert(next_idx, slide_el)

                # Add Title and Subtitle
                t_box = new_slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.1))
                t_tf = t_box.text_frame
                t_tf.word_wrap = True
                tp = t_tf.paragraphs[0]
                tp.text = "Architectural Comparison: Rolling Window vs. Checkpoint Polling"
                tp.font.name = "Calibri"
                tp.font.size = Pt(26)
                tp.font.bold = True
                tp.font.color.rgb = RGBColor(16, 44, 87)

                subp = t_tf.add_paragraph()
                subp.text = "Demystifying Historical Feature Memory Depth (W) vs. Real-Time Dashboard Execution Frequency (Δt)"
                subp.font.name = "Calibri"
                subp.font.size = Pt(15)
                subp.font.color.rgb = RGBColor(100, 110, 120)

                # Add Comparison Table: 5 rows, 3 cols
                rows = 5
                cols = 3
                t_left = Inches(0.6)
                t_top = Inches(1.6)
                t_width = Inches(12.1)
                t_height = Inches(3.9)

                t_shape = new_slide.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
                table = t_shape.table

                table.columns[0].width = Inches(2.4) # Dimension
                table.columns[1].width = Inches(4.85) # Rolling Window
                table.columns[2].width = Inches(4.85) # Polling Interval

                headers = ["Dimension / Concept", "Rolling Lookback Window (W = 10 Cycles)", "Polling Checkpoint Interval (Δt = 5 Cycles)"]
                data = [
                    ["Core Definition", "Historical memory depth looking backward in time (Cycle t minus t-10).", "Forward execution step telling the BMS microcontroller how often to run AI inference."],
                    ["Primary Function", "Feature Extraction! Calculates voltage curve variance (dQ_log_var) and capacity fade slope over past 10 cycles.", "Execution Scheduling! Controls real-time vehicle dashboard updates and system wake-ups."],
                    ["Why Chosen Value Matters", "W = 10 prevents cold-start initialization lag while filtering out instantaneous voltage sensor noise.", "Δt = 5 provides high-frequency continuous tracking to catch sudden end-of-life drops immediately."],
                    ["Everyday Driving Analogy", "The Rearview Mirror Length (How far back you look to calculate rate of change).", "The Speedometer Refresh Rate (How often the dashboard display updates for the driver)."]
                ]

                for col_i, h_text in enumerate(headers):
                    cell = table.cell(0, col_i)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(16, 44, 87)
                    p = cell.text_frame.paragraphs[0]
                    p.text = h_text
                    p.alignment = PP_ALIGN.CENTER
                    p.font.name = "Calibri"
                    p.font.size = Pt(13.5)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)

                for row_i, row_data in enumerate(data):
                    fill_col = RGBColor(245, 247, 250) if row_i % 2 == 0 else RGBColor(255, 255, 255)
                    for col_i, val in enumerate(row_data):
                        cell = table.cell(row_i + 1, col_i)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = fill_col
                        p = cell.text_frame.paragraphs[0]
                        p.text = val
                        p.alignment = PP_ALIGN.CENTER if col_i == 0 else PP_ALIGN.LEFT
                        p.font.name = "Calibri"
                        p.font.size = Pt(12.5)
                        p.font.bold = (col_i == 0)
                        p.font.color.rgb = RGBColor(33, 37, 41)

                # Add Takeaway Box below table
                box = new_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t_left, Inches(5.75), t_width, Inches(1.3))
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(235, 245, 240)
                box.line.color.rgb = RGBColor(40, 167, 69)
                box.line.width = Pt(1.5)

                btf = box.text_frame
                btf.word_wrap = True
                btf.margin_left = Inches(0.2)
                btf.margin_top = Inches(0.12)

                bp0 = btf.paragraphs[0]
                bp0.text = "Key Engineering Takeaway:"
                bp0.font.name = "Calibri"
                bp0.font.size = Pt(14)
                bp0.font.bold = True
                bp0.font.color.rgb = RGBColor(25, 135, 84)

                bp1 = btf.add_paragraph()
                bp1.text = "• Lookback Window (W) governs what historical features the AI sees from the past.\n• Polling Interval (Δt) governs how often the AI acts and warns the driver in the present.\nTogether, W = 10 and Δt = 5 form an optimal, impenetrable safety shield against EV battery failure!"
                bp1.font.name = "Calibri"
                bp1.font.size = Pt(12.5)
                bp1.font.color.rgb = RGBColor(33, 37, 41)

                print("Successfully added Comparison slide right after Slide 13.")

    try:
        prs.save(pptx_path)
        print(f"Presentation saved successfully to: {pptx_path}")
    except PermissionError:
        print(f"PERMISSION ERROR: Please close PowerPoint! File is locked: {pptx_path}")

if __name__ == "__main__":
    update_rul_final(r"reports\presentations\rul final.pptx")
    # Also save a copy to root if user prefers
    if os.path.exists(r"reports\presentations\rul final.pptx"):
        import shutil
        shutil.copy2(r"reports\presentations\rul final.pptx", r"rul final.pptx")
