import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def apply_all_modifications(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return
    
    print(f"Opening presentation: {pptx_path}")
    prs = pptx.Presentation(pptx_path)

    # -------------------------------------------------------------------------
    # 1. ADD STUDENT DETAILS TO SLIDE 1 (TITLE SLIDE)
    # -------------------------------------------------------------------------
    print("Step 1: Adding student details to Slide 1...")
    slide1 = prs.slides[0]
    
    # Remove any previously added student detail box if re-running
    for sp in list(slide1.shapes):
        if sp.has_text_frame and "R S Chandramohan" in sp.text_frame.text:
            sp._element.getparent().remove(sp._element)
            
    # Add stylish student details box at bottom left
    s_box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(5.4), Inches(6.8), Inches(1.6)
    )
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(16, 44, 87)  # Deep Navy background
    s_box.line.color.rgb = RGBColor(218, 165, 32)     # Gold accent border
    s_box.line.width = Pt(2.0)
    
    stf = s_box.text_frame
    stf.word_wrap = True
    stf.margin_left = Inches(0.2)
    stf.margin_right = Inches(0.2)
    stf.margin_top = Inches(0.15)
    
    p0 = stf.paragraphs[0]
    p0.text = "Presenter Details:"
    p0.font.name = "Calibri"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(218, 165, 32)  # Gold title
    
    details = [
        "Name: R S Chandramohan",
        "Branch: ECE  |  Year of Study: 4th Year",
        "College: Sastra Deemed to be University"
    ]
    for d in details:
        p = stf.add_paragraph()
        p.text = d
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    print(" -> Added presenter card to Slide 1.")

    # -------------------------------------------------------------------------
    # 2. REPLACE LIST ON SLIDE 6 (8 FEATURES) WITH TABLE 1 (PROPER EQUATIONS)
    # -------------------------------------------------------------------------
    print("Step 2: Transforming Slide 6 features list into Table 1 with equations...")
    slide6 = prs.slides[5]  # Index 5 is Slide 6
    
    # Remove all shapes on slide 6 except the title (Shape 0)
    for sp in list(slide6.shapes)[1:]:
        sp._element.getparent().remove(sp._element)
        
    # Table dimensions: 9 rows (1 header + 8 features), 4 cols
    rows, cols = 9, 4
    t_left, t_top = Inches(0.5), Inches(1.5)
    t_width, t_height = Inches(12.33), Inches(5.5)
    
    table_shape = slide6.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.2)   # Symbol
    table.columns[1].width = Inches(2.2)   # Feature Name
    table.columns[2].width = Inches(5.13)  # Physical Meaning & Degradation Link
    table.columns[3].width = Inches(3.8)   # Computation Formula
    
    headers = ["Symbol", "Feature Name", "Physical Meaning & Degradation Link", "Computation Formula"]
    
    # Helper to style cells and add formatted runs
    def set_cell_content(cell, runs_list, align=PP_ALIGN.LEFT, is_header=False, row_idx=0):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if is_header:
            cell.fill.fore_color.rgb = RGBColor(16, 44, 87)
        else:
            cell.fill.fore_color.rgb = RGBColor(245, 247, 250) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            
        p = cell.text_frame.paragraphs[0]
        p.text = ""
        p.alignment = align
        
        for run_data in runs_list:
            # run_data: (text, italic, bold, subscript, superscript)
            text, italic, bold, sub, sup = run_data[0], run_data[1], run_data[2], run_data[3], run_data[4]
            r = p.add_run()
            r.text = text
            r.font.name = "Calibri"
            r.font.size = Pt(13) if is_header else Pt(11.5)
            r.font.italic = italic
            r.font.bold = bold
            r.font.subscript = sub
            r.font.superscript = sup
            r.font.color.rgb = RGBColor(255, 255, 255) if is_header else RGBColor(33, 37, 41)

    # Style Header Row
    for col_i, h_text in enumerate(headers):
        set_cell_content(table.cell(0, col_i), [(h_text, False, True, False, False)], align=PP_ALIGN.CENTER, is_header=True)
        
    # Feature Data: each entry is (Symbol_runs, Name_runs, Meaning_runs, Formula_runs)
    features_data = [
        # 1. SOH
        (
            [("SOH", False, True, False, False)],
            [("State of Health", False, True, False, False)],
            [("Normalized remaining capacity ratio; primary macro-level health indicator across cell form factors.", False, False, False, False)],
            [("Q", True, True, False, False), ("dis", False, False, True, False), ("(k) / Q", True, True, False, False), ("nominal", False, False, True, False)]
        ),
        # 2. IR
        (
            [("IR", False, True, False, False)],
            [("Internal Resistance", False, True, False, False)],
            [("Measures ohmic + SEI conduction resistance; primary indicator of electrolyte decomposition and interfacial thickening.", False, False, False, False)],
            [("Mean( dV / dI )", True, True, False, False), ("10% SOC pulse", False, False, True, False)]
        ),
        # 3. dQ_log_var
        (
            [("ΔQ", True, True, False, False), ("log var", False, False, True, False)],
            [("IC Curve Log-Variance", False, True, False, False)],
            [("Captures thermodynamic peak broadening in incremental capacity curve; acts as an early warning sentinel.", False, False, False, False)],
            [("log[ Var( dQ / dV ) ]", False, True, False, False), ("W=10", False, False, True, False)]
        ),
        # 4. cycle (k)
        (
            [("k", True, True, False, False)],
            [("Cycle Number", False, True, False, False)],
            [("Tracks baseline temporal aging progression and cumulative operational charge-discharge duration of the cell.", False, False, False, False)],
            [("Current operational cycle index ", False, False, False, False), ("k", True, True, False, False), (" at checkpoint", False, False, False, False)]
        ),
        # 5. capacity_fade_window (ΔSOH)
        (
            [("ΔSOH", False, True, False, False)],
            [("Capacity Fade Rate", False, True, False, False)],
            [("Normalized rate of capacity loss over lookback window; measures local degradation velocity (ΔSOH).", False, False, False, False)],
            [("[ Q", True, True, False, False), ("dis", False, False, True, False), ("(k - 10) - Q", True, True, False, False), ("dis", False, False, True, False), ("(k) ] / Q", True, True, False, False), ("nominal", False, False, True, False)]
        ),
        # 6. Tmean
        (
            [("T", True, True, False, False), ("mean", False, False, True, False)],
            [("Mean Discharge Temp", False, True, False, False)],
            [("Monitors thermal stress during operational discharge; drives Arrhenius aging acceleration and SEI growth.", False, False, False, False)],
            [("Mean T", True, True, False, False), ("surface/core", False, False, True, False), (" over ", False, False, False, False), ("W = 10", True, True, False, False)]
        ),
        # 7. dQ_min
        (
            [("ΔQ", True, True, False, False), ("min", False, False, True, False)],
            [("IC Curve Min Shift", False, True, False, False)],
            [("Maximum localized downward shift (valley) in the differential capacity profile over the rolling lookback window.", False, False, False, False)],
            [("min[ Δ( dQ / dV ) ]", False, True, False, False), ("W=10", False, False, True, False)]
        ),
        # 8. dQ_mean
        (
            [("ΔQ", True, True, False, False), ("mean", False, False, True, False)],
            [("IC Curve Mean Shift", False, True, False, False)],
            [("Average baseline shift of incremental capacity curve; tracks global thermodynamic drift and stoichiometric imbalance.", False, False, False, False)],
            [("mean[ Δ( dQ / dV ) ]", False, True, False, False), ("W=10", False, False, True, False)]
        ),
    ]
    
    for row_idx, row_data in enumerate(features_data, start=1):
        # Symbol
        set_cell_content(table.cell(row_idx, 0), row_data[0], align=PP_ALIGN.CENTER, row_idx=row_idx)
        # Name
        set_cell_content(table.cell(row_idx, 1), row_data[1], align=PP_ALIGN.LEFT, row_idx=row_idx)
        # Meaning
        set_cell_content(table.cell(row_idx, 2), row_data[2], align=PP_ALIGN.LEFT, row_idx=row_idx)
        # Formula
        set_cell_content(table.cell(row_idx, 3), row_data[3], align=PP_ALIGN.LEFT, row_idx=row_idx)
        
    print(" -> Successfully replaced Slide 6 feature list with Table 1 equations.")

    # -------------------------------------------------------------------------
    # 3. ADD PERFORMANCE METRICS TABLE TO SLIDE 21 (CONFUSION MATRIX)
    # -------------------------------------------------------------------------
    print("Step 3: Adding Performance Metrics table to Slide 21 (Confusion Matrix)...")
    # Find Slide 21 (Confusion Matrix)
    cm_slide = None
    for slide in prs.slides:
        if slide.shapes and slide.shapes[0].has_text_frame and "Confusion Matrix" in slide.shapes[0].text_frame.text:
            cm_slide = slide
            break
            
    if cm_slide:
        # Remove any existing formula text box or card on the right side (top >= 4.0)
        for sp in list(cm_slide.shapes):
            if sp.top >= Inches(3.8) and sp.left >= Inches(6.0):
                sp._element.getparent().remove(sp._element)
                
        # Add the Performance Metrics table on the right side
        m_rows, m_cols = 5, 3
        m_left, m_top = Inches(6.4), Inches(4.2)
        m_width, m_height = Inches(6.5), Inches(3.0)
        
        m_table_shape = cm_slide.shapes.add_table(m_rows, m_cols, m_left, m_top, m_width, m_height)
        m_table = m_table_shape.table
        
        m_table.columns[0].width = Inches(1.5)  # Performance Metrics
        m_table.columns[1].width = Inches(2.3)  # Formula
        m_table.columns[2].width = Inches(2.7)  # Calculation
        
        m_headers = ["Performance Metrics", "Formula", "Calculation"]
        m_data = [
            ("Accuracy", "(TP + TN) / (TP + TN + FP + FN)  (8)", "(498 + 3600) / (498 + 3600 + 72 + 64) = 4098 / 4234 = 96.79%"),
            ("Precision", "TP / (TP + FP)  (9)", "498 / (498 + 72) = 498 / 570 = 87.37%"),
            ("Recall", "TP / (TP + FN)  (10)", "498 / (498 + 64) = 498 / 562 = 88.61%"),
            ("Overall (F1)", "F1 = (2 · Precision · Recall) / (Precision + Recall)  (11)", "(2 · 0.8737 · 0.8861) / (0.8737 + 0.8861) = 87.98%")
        ]
        
        for col_i, h_text in enumerate(m_headers):
            set_cell_content(m_table.cell(0, col_i), [(h_text, False, True, False, False)], align=PP_ALIGN.CENTER, is_header=True)
            
        for row_i, row_data in enumerate(m_data, start=1):
            set_cell_content(m_table.cell(row_i, 0), [(row_data[0], False, True, False, False)], align=PP_ALIGN.CENTER, row_idx=row_i)
            set_cell_content(m_table.cell(row_i, 1), [(row_data[1], False, False, False, False)], align=PP_ALIGN.CENTER, row_idx=row_i)
            set_cell_content(m_table.cell(row_i, 2), [(row_data[2], False, True, False, False)], align=PP_ALIGN.LEFT, row_idx=row_i)
            
        print(" -> Added Performance Metrics Table to Slide 21.")

    # -------------------------------------------------------------------------
    # 4. MAKE SLIDE 24 CONCLUSION & FUTURE WORK WITH STRUCTURED POINTS
    # -------------------------------------------------------------------------
    print("Step 4: Transforming Slide 24 into Conclusion & Future Work...")
    slide24 = prs.slides[23]  # Index 23 is Slide 24
    
    # Update Title
    title_sp = slide24.shapes[0]
    title_tf = title_sp.text_frame
    title_tf.word_wrap = True
    title_tf.paragraphs[0].text = "Conclusion & Future Work"
    title_tf.paragraphs[0].font.name = "Calibri"
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(16, 44, 87)
    
    if len(title_tf.paragraphs) > 1:
        title_tf.paragraphs[1].text = "Summary of Technical Contributions & Roadmap for Next-Generation Battery Chemistries"
        title_tf.paragraphs[1].font.name = "Calibri"
        title_tf.paragraphs[1].font.size = Pt(15)
        title_tf.paragraphs[1].font.color.rgb = RGBColor(100, 110, 120)
    else:
        sub_p = title_tf.add_paragraph()
        sub_p.text = "Summary of Technical Contributions & Roadmap for Next-Generation Battery Chemistries"
        sub_p.font.name = "Calibri"
        sub_p.font.size = Pt(15)
        sub_p.font.color.rgb = RGBColor(100, 110, 120)
        
    # Update Content Text Box (Shape 1)
    content_sp = slide24.shapes[1]
    ctf = content_sp.text_frame
    ctf.word_wrap = True
    ctf.margin_top = Inches(0.2)
    
    # Clear existing paragraphs safely
    p0 = ctf.paragraphs[0]
    p0.text = ""
    for p in list(ctf.paragraphs)[1:]:
        p._element.getparent().remove(p._element)
        
    conclusion_points = [
        ("Lightweight & Physics-Informed AI: ", "Successfully predicted LFP battery RUL without relying on flat open-circuit voltage plateaus by extracting 8 rolling-window features every 5 cycles (W=10)."),
        ("Leakage-Free Cell-Level Validation: ", "Evaluated across 124 commercial APR18650M1A cells (22,474 checkpoints) using GroupShuffleSplit, achieving robust out-of-sample generalization (MAE = 81.35 cycles, R² = 78.52%)."),
        ("Microsecond Automotive Inference: ", "Demonstrated 0.95 ms total execution time on ARM Cortex microcontrollers with an 809 KB flash footprint—leaving >99.9% CPU budget free for vehicle control operations."),
        ("Guaranteed Conformal Safety Brackets: ", "Implemented split-conformal prediction to establish a mathematically rigorous ±122 cycle safety window (90.4% coverage), eliminating black-box point forecast risks."),
        ("Electrochemical Physics Verification: ", "Bridged data science and electrochemistry via first-order Equivalent Circuit Model (ECM) transfer functions, proving that internal resistance (IR) doubling drives system poles toward instability."),
        ("Future Work – Chemistry Transfer & Adaptive Polling: ", "Extending our tree-based architecture to silicon-anode and sodium-ion (Na-ion) battery chemistries, while deploying real-time adaptive polling algorithms onto commercial automotive ECUs.")
    ]
    
    for idx, (title_text, body_text) in enumerate(conclusion_points):
        p = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
        p.space_after = Pt(8)
        
        # Title bullet
        r1 = p.add_run()
        r1.text = "• " + title_text
        r1.font.name = "Calibri"
        r1.font.size = Pt(13.5)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(16, 44, 87) if idx < 5 else RGBColor(218, 165, 32) # Gold highlight for Future Work
        
        # Body text
        r2 = p.add_run()
        r2.text = body_text
        r2.font.name = "Calibri"
        r2.font.size = Pt(13)
        r2.font.bold = False
        r2.font.color.rgb = RGBColor(33, 37, 41)
        
    print(" -> Successfully constructed Slide 24 Conclusion & Future Work.")

    # Save presentation
    try:
        prs.save(pptx_path)
        print(f"\nSUCCESS! Saved updated presentation to: {pptx_path}")
        # Also copy to root directory if present
        root_path = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(pptx_path))), "rul final.pptx")
        if os.path.exists(root_path):
            import shutil
            shutil.copy2(pptx_path, root_path)
            print(f"Synced copy to root directory: {root_path}")
    except PermissionError:
        print(f"\nPERMISSION ERROR: Please close PowerPoint! File is currently locked: {pptx_path}")

if __name__ == "__main__":
    target = r"D:\chandru project\RUL prediction\reports\presentations\rul final.pptx"
    apply_all_modifications(target)
